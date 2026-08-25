"""Real MarkItDown conversion coverage; documents are generated per test."""

from __future__ import annotations

import importlib.util
import io
from pathlib import Path
import sys
import tempfile
import unittest
import zipfile
from unittest import mock

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "src"))
import markitdown_backend as backend


MARKITDOWN_AVAILABLE = importlib.util.find_spec("markitdown") is not None
SKIP_REASON = "real conversion integration tests require the locked MarkItDown runtime"


def write_pdf(path: Path, text: str) -> None:
    """Write a tiny deterministic PDF using only its built-in Helvetica font."""
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    stream = f"BT /F1 18 Tf 72 720 Td ({text}) Tj ET".encode("ascii")
    objects.append(b"<< /Length " + str(len(stream)).encode("ascii") + b" >>\nstream\n" + stream + b"\nendstream")
    data = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for number, body in enumerate(objects, 1):
        offsets.append(len(data))
        data.extend(f"{number} 0 obj\n".encode("ascii"))
        data.extend(body)
        data.extend(b"\nendobj\n")
    xref = len(data)
    data.extend(f"xref\n0 {len(objects) + 1}\n0000000000 65535 f \n".encode("ascii"))
    data.extend(b"".join(f"{offset:010d} 00000 n \n".encode("ascii") for offset in offsets[1:]))
    data.extend(f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF\n".encode("ascii"))
    path.write_bytes(data)


def write_docx(path: Path, text: str) -> None:
    """Write the minimal OOXML package MarkItDown needs, without a test dependency."""
    files = {
        "[Content_Types].xml": """<?xml version=\"1.0\" encoding=\"UTF-8\"?>
<Types xmlns=\"http://schemas.openxmlformats.org/package/2006/content-types\">
  <Default Extension=\"rels\" ContentType=\"application/vnd.openxmlformats-package.relationships+xml\"/>
  <Default Extension=\"xml\" ContentType=\"application/xml\"/>
  <Override PartName=\"/word/document.xml\" ContentType=\"application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml\"/>
</Types>""",
        "_rels/.rels": """<?xml version=\"1.0\" encoding=\"UTF-8\"?>
<Relationships xmlns=\"http://schemas.openxmlformats.org/package/2006/relationships\">
  <Relationship Id=\"rId1\" Type=\"http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument\" Target=\"word/document.xml\"/>
</Relationships>""",
        "word/document.xml": f"""<?xml version=\"1.0\" encoding=\"UTF-8\"?>
<w:document xmlns:w=\"http://schemas.openxmlformats.org/wordprocessingml/2006/main\">
  <w:body><w:p><w:r><w:t>{text}</w:t></w:r></w:p></w:body>
</w:document>""",
    }
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for name, contents in files.items():
            archive.writestr(name, contents)


def write_pptx(path: Path, text: str) -> None:
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = text
    presentation.save(path)


def write_xlsx(path: Path, text: str) -> None:
    from openpyxl import Workbook

    workbook = Workbook()
    workbook.active["A1"] = text
    workbook.save(path)


@unittest.skipUnless(MARKITDOWN_AVAILABLE, SKIP_REASON)
class RealConversionTests(unittest.TestCase):
    def convert(self, suffix: str, writer, *, difficult_name: bool = False, repeat: bool = False) -> None:
        token = f"REAL-{suffix[1:].upper()}-FIXTURE"
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            name = f"report {suffix}" if not difficult_name else f"report ü $(); 'quoted'\nline {suffix}"
            source = root / name
            writer(source, token)
            clipboard_payloads: list[str] = []
            notifications: list[tuple[str, str, str]] = []

            def clipboard(text: str) -> tuple[bool, str]:
                clipboard_payloads.append(text)
                return True, ""

            def notifier(headline: str, description: str, urgency: str) -> bool:
                notifications.append((headline, description, urgency))
                return True

            with mock.patch("sys.stdout", new_callable=io.StringIO) as output:
                code = backend.run_conversion(
                    str(source), destination_dir=root, clipboard=clipboard, notifier=notifier
                )
            expected = root / f"{source.stem}.md"
            self.assertEqual(code, 0)
            self.assertTrue(expected.is_file())
            markdown = expected.read_text(encoding="utf-8")
            self.assertTrue(markdown.strip())
            self.assertIn(token, markdown)
            self.assertEqual(clipboard_payloads, [markdown])
            self.assertEqual(len(notifications), 1)
            self.assertIn('"status": "success"', output.getvalue())
            self.assertFalse(list(root.glob(".markitdown-source-*")))
            self.assertFalse(list(root.glob(".markitdown-*.tmp")))

            if repeat:
                with mock.patch("sys.stdout", new_callable=io.StringIO):
                    code = backend.run_conversion(
                        str(source), destination_dir=root, clipboard=clipboard, notifier=notifier
                    )
                collision = root / f"{source.stem}-2.md"
                self.assertEqual(code, 0)
                self.assertTrue(collision.is_file())
                self.assertIn(token, collision.read_text(encoding="utf-8"))
                self.assertFalse(list(root.glob(".markitdown-source-*")))
                self.assertFalse(list(root.glob(".markitdown-*.tmp")))

    def test_pdf_with_difficult_name_and_collision(self) -> None:
        self.convert(".pdf", write_pdf, difficult_name=True, repeat=True)

    def test_docx(self) -> None:
        self.convert(".docx", write_docx)

    def test_pptx(self) -> None:
        self.convert(".pptx", write_pptx)

    def test_xlsx(self) -> None:
        self.convert(".xlsx", write_xlsx)


if __name__ == "__main__":
    unittest.main()
